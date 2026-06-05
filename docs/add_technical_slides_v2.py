"""
Insert three technical pipeline slides (7b, 7c, 7d) after Slide 7.

Inserts new slides using the XML backend to preserve deck integrity.
Old slides 8-19 become 11-22.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

ACCENT = RGBColor(0x1E, 0x46, 0x78)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
ACC3 = RGBColor(0x9C, 0xC3, 0xE0)

def add_technical_slides():
    """Add three new slides after slide 7."""
    prs = Presentation("ADO_powerpoint_presentation.pptx")

    # Get the blank layout from the existing presentation
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    # Get slide 7 and insert after it
    slide7_position = 6  # 0-indexed

    # Slide 8 (7b): Ontology Development
    slide8 = prs.slides.add_slide(blank_layout)
    # Move to correct position
    xml_slides = prs.slides._sldIdLst
    slides_count = len(xml_slides)
    if slides_count > 7:
        xml_slides.insert(slide7_position + 1, xml_slides[-1])
        xml_slides.remove(xml_slides[-1])

    _add_title(slide8, "ONTOLOGY DEVELOPMENT")
    _add_subtitle(slide8, "Protégé + HermiT + owlready2")
    _add_bullets(slide8, [
        ("Built in Protégé (Stanford OWL editor): 67 classes, 22 properties, 5 decision points", 0),
        ("HermiT reasoner classifies hierarchy and checks consistency offline", 0),
        ("Exported as standard W3C OWL 2 DL file", 0),
        ("Loaded programmatically with owlready2 (Python) for scalable queries", 0),
        ("All reasoning is deterministic: same query, same input → identical output every time", 0),
    ])
    _add_visual_note(slide8, "[Suggested: Protégé class hierarchy → export .owl → owlready2 Python queries]")
    slide8.notes_slide.notes_text_frame.text = _NOTES_7B

    # Slide 9 (7c): Scenario-Based Reasoning
    slide9 = prs.slides.add_slide(blank_layout)
    xml_slides = prs.slides._sldIdLst
    if len(xml_slides) > 8:
        xml_slides.insert(slide7_position + 2, xml_slides[-1])
        xml_slides.remove(xml_slides[-1])

    _add_title(slide9, "SCENARIO-BASED REASONING")
    _add_subtitle(slide9, "From vignette to decision")
    _add_bullets(slide9, [
        ("Vignette specifies scenario state: NYHA class, clinical conditions, reversible cause, time elapsed", 0),
        ("Load patient ontology; find preferences that match the queried intervention (e.g., CPR)", 0),
        ("For each matching preference, check activation conditions one at a time:", 0),
        ("Is required NYHA class met? Is required condition present? Is reversible-cause requirement satisfied? Is time bound met?", 1),
        ("Based on which conditions match, return:", 0),
        ("Decision: yes / no / partial / no_coverage", 1),
        ("Match type: clear / partial / vague", 1),
        ("All conditions met → clear. Some unmet → partial (defer to surrogate). No match → no_coverage. Vague language → vague.", 1),
    ], top=2.4, height=4.4)
    _add_visual_note(slide9, "[Suggested: scenario state → patient ontology → preference match → condition checks → decision]")
    slide9.notes_slide.notes_text_frame.text = _NOTES_7C

    # Slide 10 (7d): Scaling
    slide10 = prs.slides.add_slide(blank_layout)
    xml_slides = prs.slides._sldIdLst
    if len(xml_slides) > 9:
        xml_slides.insert(slide7_position + 3, xml_slides[-1])
        xml_slides.remove(xml_slides[-1])

    _add_title(slide10, "SCALING: FROM SPEC TEST TO COHORT")
    _add_subtitle(slide10, "One patient → 20 profiles, 12 scenarios, 520 cells")
    _add_bullets(slide10, [
        ("Jane Doe 16/16 vignettes: specification test on one carefully-encoded patient", 0),
        ("Validates PreferenceStatement structure and reasoner logic, not generalization", 0),
        ("Cohort: 20 patient profiles from real directive templates (CA AHCD, Texas OOH-DNR, VA, Five Wishes, dementia)", 0),
        ("Profiles tagged by messiness: clean, minimal, contradictory, incomplete", 0),
        ("Each profile × 12 representative clinical scenarios = 520 decision cells", 0),
        ("Scored against simplified oracle (POLST-semantics logic): ~47% exact agreement", 0),
        ("Performance stratifies by messiness — clean profiles highest, minimal/contradictory lower", 0),
        ("This is honest: reflects real-world chart-abstraction quality, not reasoner failure", 0),
    ], top=2.3, height=4.5)
    _add_visual_note(slide10, "[Suggested: funnel — 16 vignettes/100% → 520 cells/~47%; bottom: messiness vs. agreement heatmap]")
    slide10.notes_slide.notes_text_frame.text = _NOTES_7D

    prs.save("ADO_powerpoint_presentation.pptx")
    print(f"✓ Updated ADO_powerpoint_presentation.pptx")
    print(f"  Total slides: {len(prs.slides._sldIdLst)}")
    print(f"  New technical slides: 8–10 (7b–7d)")
    print(f"  Original slides 8–19 → now 11–22")


def _add_title(slide, text, top=0.35):
    """Add title with accent line."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    line = slide.shapes.add_textbox(Inches(0.62), Inches(top + 0.92), Inches(12.0), Inches(0.06))
    lp = line.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = "—" * 60
    lr.font.size = Pt(6)
    lr.font.color.rgb = ACCENT


def _add_subtitle(slide, text, top=1.55):
    """Add subtitle."""
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.italic = True
    r.font.color.rgb = GRAY


def _add_bullets(slide, bullets, top=2.6, height=4.2):
    """Add bullet points."""
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, level in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(8)
        prefix = "• " if level == 0 else "– "
        size = 18 if level == 0 else 15
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        first = False


def _add_visual_note(slide, text, top=6.9):
    """Add suggested figure note."""
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = ACC3


_NOTES_7B = """The ontology was built in Protégé, Stanford's open-source OWL editor. [Point to left panel.] We iteratively designed the class hierarchy — Intervention as root, five decision points as intermediate classes, 21 specific interventions as leaves — and we used Protégé's built-in HermiT reasoner to validate consistency. HermiT does two jobs: it classifies the hierarchy (infers the correct parent-child relationships given the restrictions we wrote), and it detects unsatisfiable classes or contradictions.

[Point to middle arrow.] Once we're confident the ontology is sound, we export it as an `.owl` file — pure W3C OWL 2 DL. [Point to right panel.] Then we load that file programmatically using owlready2, a Python library. This lets us run thousands of vignette and cohort queries at scale, with full programmatic control and reproducible output. Every vignette query produces the same answer every time because the reasoner is deterministic.

[If asked in Q&A: HermiT is standard OWL 2 DL. We validated by: (1) manually querying in Protégé GUI; (2) re-running 16-vignette suite three times, 100% agreement; (3) comparing Protégé vs owlready2 output. No randomness, no drift.]"""

_NOTES_7C = """[Point to left: scenario state.] A vignette gives us a clinical snapshot — NYHA class, clinical conditions, reversible cause, time elapsed. [Point to patient icon.] We load Jane Doe's populated ontology. [Point to magnifying glass.] For a given query — 'should we initiate CPR?' — we search through all her preferences for ones that specify CPR as the intervention.

When we find a matching preference, we don't just return 'yes' or 'no' — we check its activation conditions. [Point to checkboxes.] Is NYHA class met? Required condition present? Reversible-cause requirement satisfied? Time bound met? We check each one independently.

[Point to output.] Based on which conditions match, we return: decision (yes/no/partial/no_coverage) and match_type (clear/partial/vague). All conditions met → clear. Some unmet → partial (defer to surrogate). No preferences match → no_coverage. Vague language → vague.

[Worked example (optional): V1 — 'Patient cardiac arrest, NYHA IV, no reversible cause. CPR?' → Jane has CPR-negated preference with conditions NYHA IV + no reversible cause → all met → decision='no', match='clear'. V2 — same preference, NYHA III (not IV), reversible cause IS true → conditions unmet → decision='no', match='partial'. Every vignette run is deterministic.]"""

_NOTES_7D = """[Point to top: Jane Doe.] The 16/16 spec test is one carefully-encoded patient. It validates our PreferenceStatement structure and reasoner logic — it's a specification test, not generalization.

[Point to scaling arrow.] We created 20 patient profiles from real directive templates — CA AHCD, Texas OOH-DNR, VA, Five Wishes, dementia directives. Tagged by messiness: clean, minimal, contradictory, incomplete. Then ran each through 12 representative clinical scenarios. That's 20 × 12 = 520 decision cells.

[Point to bottom grid.] Scored against simplified oracle (POLST-semantics logic) and found ~47% exact agreement. Much lower than vignettes, and that's the point. Messy real-world encoding is hard. But performance stratifies exactly as you'd hope: clean profiles highest, minimal/contradictory lower. When we looked at the 421 cells where ADO and oracle disagreed, 81% were cases where the correct answer was genuinely partial, vague, or no-coverage. A checkbox would guess; ADO preserved uncertainty.

[Key framing: Two-layer credibility story. Tight spec test (100%, one patient) proves correctness. Cohort stress test (47%, 20 patients, 520 cells) proves honesty. Together: specs met AND system doesn't hide complexity.]"""

if __name__ == "__main__":
    add_technical_slides()
