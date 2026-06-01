#!/usr/bin/env python3
"""Build docs/presentation/ADO_powerpoint_presentation.pptx from the slide outline in Presentation.tex."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ADO_BLUE = RGBColor(30, 70, 120)
OUTPUT = Path(__file__).resolve().parent / "docs/presentation/ADO_powerpoint_presentation.pptx"


def add_bullets(text_frame, lines, level=0, bold_first=False):
    text_frame.clear()
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.level = level
        p.font.size = Pt(18 if level == 0 else 16)
        if i == 0 and bold_first:
            p.font.bold = True


def add_title_slide(prs, title, subtitle, authors, institute, date):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"{subtitle}\n\n{authors}\n{institute}\n{date}"


def add_content_slide(prs, title, bullets, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    add_bullets(body, bullets)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_two_column_slide(prs, title, left_title, left_bullets, right_block_title, right_block_text):
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # two content
    slide.shapes.title.text = title
    left = slide.placeholders[1].text_frame
    add_bullets(left, [left_title] + left_bullets, bold_first=True)
    right = slide.placeholders[2].text_frame
    add_bullets(right, [right_block_title, "", right_block_text])


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "The Advance Directive Ontology (ADO)",
        "A Computable Representation of End-of-Life Care Preferences in Advanced Heart Failure",
        "Patrick Walsh and Darren Chan",
        "Stanford University | BMDS 210 / CS 270",
        "June 1, 2026",
    )

    add_content_slide(
        prs,
        "The problem: directives that fail at the bedside",
        [
            "67% of ICU patients lack decision-making capacity at end of life; fewer than 37% of US adults have any directive on file.",
            "The usual story — directives are unreadable free text — is only partly true: good forms (California, UC) are mostly checkboxes.",
            "REAL PROBLEM: existing instruments are low-resolution and disease-blind:",
            "  • They collapse conditional, graded preferences into a few boxes",
            "  • They ignore HF-specific decisions: ICD, LVAD, inotropes",
            "  • At 2 a.m. it is POLST and code-status orders, not the directive, that get acted on",
        ],
    )

    add_two_column_slide(
        prs,
        "What we built",
        "ADO — an OWL 2 ontology of EOL preferences:",
        [
            "67 classes, 22 properties; 5 HF decision points; 21 interventions",
            "Encodes conditional, graded, negatable preferences + verbatim language",
            "Reasoner returns one of four honest outputs",
        ],
        "Our stance",
        "ADO is a decision-support aid that helps pre-populate an advance-care-planning note. It is NOT a legal directive — and complements POLST and the code-status system.",
    )

    add_content_slide(
        prs,
        "Architecture: ontology with LLMs at the boundaries",
        [
            "Free-text directive → LLM extraction → OWL ontology → Reasoner → Decision + code status",
            "LLM extracts (closed-world: encode only what the patient stated)",
            "Ontology is the auditable, closed-world source of truth",
            "Reasoner produces the defensible four-category output",
            "",
            "A pure LLM is non-deterministic and over-confident; an auditable reasoner can honestly say the directive is silent.",
        ],
    )

    add_content_slide(
        prs,
        "The honest four-category output",
        [
            "Clear — activation conditions met; preference unambiguous",
            "Partial — some conditions unmet → defer to surrogate",
            "No coverage — the directive is silent → escalate, do not presume",
            "Vague — surface the patient's own words (e.g. no heroic measures)",
            "",
            "A checkbox — or a confident LLM — cannot say I don't know. ADO can, and that is the point.",
        ],
    )

    add_two_column_slide(
        prs,
        "Results I: reasoning + extraction",
        "Scenario-based reasoning",
        [
            "16 clinical vignettes, all 5 decision points, all 4 preference types",
            "16/16 decision and match-type accuracy",
        ],
        "LLM extraction (real directive text)",
        "Precision 0.94 / Recall 1.00 / F1 0.97. Negation, type, conditionality: 15/15. Zero hallucinations on out-of-scope text.",
    )

    add_content_slide(
        prs,
        "Results II (flagship): can ADO produce what governs the bedside?",
        [
            "12 preference profiles from real templates (CA AHCD, Texas OOH-DNR, NHS Wales ADRT, Five Wishes, VA, dementia directive)",
            "Gold standard = POLST semantics (independent of ADO)",
            "Hospital code status: 12/12 | POLST Section A: 12/12 | POLST Section B: 11/12",
            "Overall: 35/36 (97%) | Cohen's κ = 1.00 / 1.00 / 0.88",
            "Value-add: same directive, two scenarios — DNR ↔ Full code as conditions are met or not",
        ],
    )

    add_content_slide(
        prs,
        "Expert review: Dr. David Magnus (Stanford Biomedical Ethics)",
        [
            "Reframed: You can't think of these as AHCDs… This is more like a progress note.",
            "On premise: It isn't largely a problem of free text — the good forms are checkboxes.",
            "Deepest critique: preferences are two-dimensional — what they want AND who decides.",
            "Response: repositioned ADO as an ACP-note aid; built code-status / POLST output; who-decides axis is next.",
        ],
    )

    add_content_slide(
        prs,
        "Limitations & future work",
        [
            "No legal force — ADO is one structured input, not an oracle",
            "One-dimensional today — add who-decides / surrogate-override axis (~54% would let family override)",
            "Add artificial hydration/nutrition and POLST time-limited trial",
            "Validate against real EHR code-status data (MIMIC-IV) — credentialing in progress",
        ],
    )

    add_content_slide(
        prs,
        "Conclusion",
        [
            "ADO is the missing inference layer above POLST and the code-status system",
            "Clinically honest four-category output — including honest non-answers",
            "16/16 reasoning, F1 0.97 extraction (zero out-of-scope hallucinations), 97% code-status mapping",
            "Ontology with LLMs at the boundaries — not ontology versus LLM",
            "",
            "Thank you — questions?",
            "https://github.com/patrickwalsh26/bmds210-AD_Ontology",
        ],
    )

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
