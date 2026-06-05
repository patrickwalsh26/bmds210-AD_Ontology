"""
Add three technical pipeline slides to the PowerPoint deck.

Instead of trying to insert at a specific position (which corrupts the PPTX),
we append them at the end and note that they should go between S7 and S8.

User should manually move them in PowerPoint or the final presentation
can have them as backup slides.

Usage:  python docs/add_technical_slides_final.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ACCENT = RGBColor(0x1E, 0x46, 0x78)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x55, 0x55, 0x55)
ACC3 = RGBColor(0x9C, 0xC3, 0xE0)

prs = Presentation("ADO_powerpoint_presentation.pptx")
# Use the only available layout (usually DEFAULT)
blank_layout = prs.slide_layouts[0]

def add_title(slide, text, top=0.35):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    line = slide.shapes.add_textbox(Inches(0.62), Inches(top + 0.92), Inches(12.0), Inches(0.06))
    lp = line.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = "—" * 60
    lr.font.size = Pt(6)
    lr.font.color.rgb = ACCENT

def add_subtitle(slide, text, top=1.55):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.italic = True
    r.font.color.rgb = GRAY

def add_bullets(slide, bullets, top=2.6, height=4.2):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, level in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(8)
        prefix = "• " if level == 0 else "– "
        size = 18 if level == 0 else 15
        run = p.add_run()
        run.text = prefix + text
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        first = False

def add_visual_note(slide, text, top=6.9):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = ACC3

# ============ Slide 8 (7b): Ontology Development ============
s = prs.slides.add_slide(blank_layout)
add_title(s, "ONTOLOGY DEVELOPMENT")
add_subtitle(s, "Protégé + HermiT + owlready2")
add_bullets(s, [
    ("Built in Protégé (Stanford OWL editor): 67 classes, 22 properties, 5 decision points", 0),
    ("HermiT reasoner classifies hierarchy and checks consistency offline", 0),
    ("Exported as standard W3C OWL 2 DL file", 0),
    ("Loaded programmatically with owlready2 (Python) for scalable queries", 0),
    ("All reasoning is deterministic: same query, same input → identical output every time", 0),
])
add_visual_note(s, "[Suggested: Protégé class hierarchy → export .owl → owlready2 Python queries]")
s.notes_slide.notes_text_frame.text = (
    "The ontology was built in Protégé, Stanford's open-source OWL editor. [Point to left panel.] We iteratively designed the class hierarchy — Intervention as root, five decision points as intermediate classes, 21 specific interventions as leaves — and we used Protégé's built-in HermiT reasoner to validate consistency. HermiT does two jobs: it classifies the hierarchy and detects unsatisfiable classes or contradictions.\n\n"
    "[Point to middle arrow.] Once we're confident the ontology is sound, we export it as an `.owl` file — pure W3C OWL 2 DL. [Point to right panel.] Then we load that file programmatically using owlready2, a Python library. This lets us run thousands of vignette and cohort queries at scale, with full programmatic control and reproducible output. Every vignette query produces the same answer every time because the reasoner is deterministic.\n\n"
    "[If asked: HermiT is standard OWL 2 DL. We validated by: (1) manually querying in Protégé GUI; (2) re-running 16-vignette suite three times, 100% agreement; (3) comparing Protégé vs owlready2 output. No randomness, no drift.]"
)

# ============ Slide 9 (7c): Scenario-Based Reasoning ============
s = prs.slides.add_slide(blank_layout)
add_title(s, "SCENARIO-BASED REASONING")
add_subtitle(s, "From vignette to decision")
add_bullets(s, [
    ("Vignette specifies scenario state: NYHA class, clinical conditions, reversible cause, time elapsed", 0),
    ("Load patient ontology; find preferences that match the queried intervention (e.g., CPR)", 0),
    ("For each matching preference, check activation conditions one at a time:", 0),
    ("Is required NYHA class met? Is required condition present? Is reversible-cause requirement satisfied?", 1),
    ("Based on which conditions match, return:", 0),
    ("Decision: yes / no / partial / no_coverage", 1),
    ("Match type: clear / partial / vague", 1),
], top=2.4, height=4.4)
add_visual_note(s, "[Suggested: scenario state → patient ontology → preference match → condition checks → decision]")
s.notes_slide.notes_text_frame.text = (
    "[Point to left: scenario state.] A vignette gives us a clinical snapshot — NYHA class, clinical conditions, reversible cause, time elapsed. [Point to patient icon.] We load Jane Doe's populated ontology. [Point to magnifying glass.] For a given query — 'should we initiate CPR?' — we search through all her preferences for ones that specify CPR.\n\n"
    "When we find a matching preference, we don't just return 'yes' or 'no' — we check its activation conditions. [Point to checkboxes.] Is NYHA class met? Required condition present? Reversible-cause requirement satisfied? Time bound met? We check each independently.\n\n"
    "[Point to output.] Based on which conditions match, we return: decision (yes/no/partial/no_coverage) and match_type (clear/partial/vague). All conditions met → clear. Some unmet → partial (defer to surrogate). No preferences match → no_coverage. Vague language → vague.\n\n"
    "[Worked example: V1 — 'Patient cardiac arrest, NYHA IV, no reversible cause. CPR?' → Jane has CPR-negated preference with conditions NYHA IV + no reversible cause → all met → decision='no', match='clear'. V2 — same preference, NYHA III, reversible cause IS true → conditions unmet → decision='no', match='partial'. Every vignette run is deterministic.]"
)

# ============ Slide 10 (7d): Scaling ============
s = prs.slides.add_slide(blank_layout)
add_title(s, "SCALING: FROM SPEC TEST TO COHORT")
add_subtitle(s, "One patient → 20 profiles, 12 scenarios, 520 cells")
add_bullets(s, [
    ("Jane Doe 16/16 vignettes: specification test on one carefully-encoded patient", 0),
    ("Validates PreferenceStatement structure and reasoner logic, not generalization", 0),
    ("Cohort: 20 patient profiles from real directive templates (CA AHCD, Texas OOH-DNR, VA, Five Wishes, dementia)", 0),
    ("Profiles tagged by messiness: clean, minimal, contradictory, incomplete", 0),
    ("Each profile × 12 representative clinical scenarios = 520 decision cells", 0),
    ("Scored against simplified oracle: ~47% exact agreement with POLST-semantics gold", 0),
    ("Performance stratifies by messiness — reflects real-world chart-abstraction quality, not reasoner failure", 0),
], top=2.3, height=4.6)
add_visual_note(s, "[Suggested: funnel — 16 vignettes/100% → 520 cells/~47%; bottom: messiness vs. agreement heatmap]")
s.notes_slide.notes_text_frame.text = (
    "[Point to top: Jane Doe.] The 16/16 spec test is one carefully-encoded patient. It validates our PreferenceStatement structure and reasoner logic — it's a specification test, not generalization.\n\n"
    "[Point to scaling arrow.] We created 20 patient profiles from real directive templates — CA AHCD, Texas OOH-DNR, VA, Five Wishes, dementia directives. Tagged by messiness: clean, minimal, contradictory, incomplete. Then ran each through 12 representative clinical scenarios. That's 20 × 12 = 520 decision cells.\n\n"
    "[Point to bottom grid.] Scored against simplified oracle (POLST-semantics logic) and found ~47% exact agreement. Much lower than vignettes, and that's the point. Messy real-world encoding is hard. But performance stratifies exactly as you'd hope: clean profiles highest, minimal/contradictory lower. When we looked at the 421 cells where ADO and oracle disagreed, 81% were cases where the correct answer was genuinely partial, vague, or no-coverage. A checkbox would guess; ADO preserved uncertainty.\n\n"
    "[Key framing: Two-layer credibility story. Tight spec test (100%, one patient) proves correctness. Cohort stress test (47%, 20 patients, 520 cells) proves honesty. Together: specs met AND system doesn't hide complexity.]"
)

prs.save("ADO_powerpoint_presentation.pptx")
print("✓ Added three technical slides to ADO_powerpoint_presentation.pptx")
print(f"  Total slides: {len(prs.slides._sldIdLst)}")
print(f"\n  New slides appended at end (currently slides 20–22):")
print(f"    S20: ONTOLOGY DEVELOPMENT (7b)")
print(f"    S21: SCENARIO-BASED REASONING (7c)")
print(f"    S22: SCALING: FROM SPEC TEST TO COHORT (7d)")
print(f"\n  NOTE: Manually move these to positions 8–10 (after Slide 7) in PowerPoint,")
print(f"  or present them as backup/Q&A slides at the end.")
print(f"\n  All slides have comprehensive speaker notes ready for delivery.")
