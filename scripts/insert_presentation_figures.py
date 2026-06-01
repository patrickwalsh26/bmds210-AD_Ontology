#!/usr/bin/env python3
"""Insert publication-quality figures and captions into docs/presentation/ADO_powerpoint_presentation.pptx."""

import sys
from pathlib import Path as _Path
_ADO_ROOT = _Path(__file__).resolve().parents[1]
if str(_ADO_ROOT) not in sys.path:
    sys.path.insert(0, str(_ADO_ROOT))

import sys
from pathlib import Path
_ROOT = Path(__file__).resolve().parents[1]
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs/presentation/ADO_powerpoint_presentation.pptx"
FIG = ROOT / "docs" / "presentation" / "presentation_figures"
CAPTIONS_JSON = FIG / "figure_captions.json"


def load_short_captions() -> dict[str, str]:
    """First sentence of each figure caption (slide-friendly)."""
    if not CAPTIONS_JSON.exists():
        return {}
    raw = json.loads(CAPTIONS_JSON.read_text(encoding="utf-8"))
    out = {}
    for name, text in raw.items():
        clean = re.sub(r"\*+", "", text)
        sentence = clean.split(".")[0].strip() + "." if clean.strip() else ""
        out[name] = sentence
    return out


def insert_on_slide(slide, img_path: Path, left, top, width):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), left, top, width=width)


def add_caption_box(slide, text: str, top=Inches(6.35), height=Inches(0.65)):
    if not text:
        return
    box = slide.shapes.add_textbox(Inches(0.45), top, Inches(12.4), height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = None  # default


def slide_has_title(prs, needle: str) -> bool:
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and needle in sh.text_frame.text:
                return True
    return False


def enhance_slide_5(prs: Presentation) -> None:
    slide = prs.slides[4]
    insert_on_slide(slide, FIG / "ontology_class_hierarchy.png", Inches(6.9), Inches(1.15), Inches(5.9))


def enhance_slide_8(prs: Presentation) -> None:
    slide = prs.slides[7]
    insert_on_slide(slide, FIG / "study_design_strands.png", Inches(0.4), Inches(1.0), Inches(12.4))
    caps = load_short_captions()
    add_caption_box(slide, caps.get("study_design_strands.png", ""), top=Inches(6.5))


def enhance_slide_9(prs: Presentation) -> None:
    slide = prs.slides[8]
    insert_on_slide(slide, FIG / "eval_two_layer.png", Inches(0.35), Inches(0.95), Inches(12.5))
    caps = load_short_captions()
    add_caption_box(slide, caps.get("eval_two_layer.png", ""), top=Inches(6.45))


def enhance_slide_10(prs: Presentation) -> None:
    slide = prs.slides[9]
    caps = load_short_captions()
    insert_on_slide(slide, FIG / "ablation_conditions.png", Inches(0.35), Inches(1.2), Inches(5.6))
    insert_on_slide(slide, FIG / "conditional_p9_p10.png", Inches(6.75), Inches(1.35), Inches(6.0))
    add_caption_box(
        slide,
        caps.get("ablation_conditions.png", "") + " " + caps.get("conditional_p9_p10.png", ""),
        top=Inches(6.35),
        height=Inches(0.75),
    )


def add_cohort_stress_slide(prs: Presentation) -> None:
    if slide_has_title(prs, "COHORT STRESS TEST"):
        return
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    caps = load_short_captions()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.55))
    tb.text_frame.text = "COHORT STRESS TEST — 20 profiles × 12 scenarios (520 cells)"
    tb.text_frame.paragraphs[0].font.size = Pt(17)
    tb.text_frame.paragraphs[0].font.bold = True

    insert_on_slide(slide, FIG / "cohort_messy_breakdown.png", Inches(0.35), Inches(0.9), Inches(6.15))
    insert_on_slide(slide, FIG / "cohort_baseline_comparison.png", Inches(6.65), Inches(0.9), Inches(6.15))
    add_caption_box(
        slide,
        caps.get("cohort_messy_breakdown.png", "") + " " + caps.get("cohort_baseline_comparison.png", ""),
        top=Inches(6.55),
        height=Inches(0.7),
    )


def add_eval_dashboard_slide(prs: Presentation) -> None:
    title = "EVALUATION DASHBOARD"
    if slide_has_title(prs, title):
        return
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    caps = load_short_captions()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5))
    tb.text_frame.text = f"{title} — developmental validation (June 2026)"
    tb.text_frame.paragraphs[0].font.size = Pt(17)
    tb.text_frame.paragraphs[0].font.bold = True
    insert_on_slide(slide, FIG / "eval_dashboard.png", Inches(0.3), Inches(0.72), Inches(12.6))
    add_caption_box(slide, caps.get("eval_dashboard.png", ""), top=Inches(6.55))


def add_coverage_extraction_slide(prs: Presentation) -> None:
    if slide_has_title(prs, "COVERAGE & EXTRACTION"):
        return
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    caps = load_short_captions()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.5))
    tb.text_frame.text = "COVERAGE & EXTRACTION — scope and closed-world discipline"
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.size = Pt(17)

    insert_on_slide(slide, FIG / "coverage_inventory.png", Inches(0.35), Inches(0.9), Inches(6.2))
    insert_on_slide(slide, FIG / "extraction_f1.png", Inches(6.75), Inches(1.0), Inches(5.7))
    add_caption_box(
        slide,
        caps.get("coverage_inventory.png", "") + " " + caps.get("extraction_f1.png", ""),
        top=Inches(6.5),
        height=Inches(0.7),
    )


def add_track3_slide(prs: Presentation) -> None:
    if slide_has_title(prs, "TRACK 3 — CODE STATUS"):
        return
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    caps = load_short_captions()
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12), Inches(0.5))
    title_box.text_frame.text = "TRACK 3 — CODE STATUS & POLST MAPPING"
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.size = Pt(17)
    insert_on_slide(slide, FIG / "track3_field_agreement.png", Inches(0.35), Inches(0.9), Inches(6.1))
    insert_on_slide(slide, FIG / "code_status_confusion.png", Inches(6.65), Inches(0.9), Inches(6.1))
    add_caption_box(
        slide,
        caps.get("track3_field_agreement.png", "") + " " + caps.get("code_status_confusion.png", ""),
        top=Inches(6.5),
        height=Inches(0.7),
    )


def main():
    if not PPTX.exists():
        raise SystemExit(f"Missing {PPTX}")
    prs = Presentation(str(PPTX))
    enhance_slide_5(prs)
    enhance_slide_8(prs)
    enhance_slide_9(prs)
    enhance_slide_10(prs)
    add_cohort_stress_slide(prs)
    add_eval_dashboard_slide(prs)
    add_coverage_extraction_slide(prs)
    add_track3_slide(prs)
    prs.save(str(PPTX))
    print(f"Updated {PPTX} ({len(prs.slides)} slides)")
    print("Captions: docs/presentation/presentation_figures/FIGURE_CAPTIONS.md")


if __name__ == "__main__":
    main()
