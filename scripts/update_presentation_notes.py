#!/usr/bin/env python3
"""Embed speaker notes and refresh slide 9 metrics in ADO_powerpoint_presentation.pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "ADO_powerpoint_presentation.pptx"

NOTES = [
    # Slide 1
    """[~30s] Good morning. Patrick Walsh and Darren Chan, Stanford BMDS 210 / CS 270.

Hook: ADO is a computable OWL 2 representation of end-of-life preferences in advanced heart failure—not a legal advance directive, but a decision-support layer that helps clinicians reason over what patients said and pre-populate ACP documentation.

Transition → why the bedside workflow breaks down.""",

    # Slide 2
    """[~45s] Three instruments, three jobs: advance directive (values/wishes), POLST (portable orders), code status (what we act on at 2 a.m.).

The loss is in the middle—conditional wishes become coarse checkboxes or bedside orders.

Important reframe (Magnus): the problem is NOT mainly “unreadable free text.” CA and UC forms are mostly checkboxes. They are low-resolution and disease-blind.

Point to the ventilator example: temporary-if-reversible vs never-indefinitely cannot fit a yes/no box.

Transition → why we scoped to advanced HF.""",

    # Slide 3
    """[~40s] HF is unpredictable, device-dependent, and patients revisit the ED/ICU—preference interpretation is time-pressured and repeated.

Five decision points on slide: CPR, ventilation, devices, vasopressors/inotropes, dialysis.

~200K HF patients have ICDs; generic ADs rarely mention ICD/LVAD/inotropes—that gap motivates our scope.

Transition → how we represent a preference computationally.""",

    # Slide 4
    """[~45s] Core unit: PreferenceStatement = intervention + activation conditions + strength + negation + verbatim originalText.

Walk the vasopressor example: 72h “no improvement” becomes a structured activation condition the reasoner can test.

Why verbatim text? Auditable—clinicians see patient language behind the inference.

Transition → ontology size and granularity.""",

    # Slide 5
    """[~40s] 67 classes, 22 properties, 21 intervention subclasses—finer than any single template in our 50-form inventory (e.g., six ventilation types vs one checkbox).

SNOMED CT grounding where possible for interoperability.

Transition → we don’t force binary answers.""",

    # Slide 6
    """[~45s] Four honest outputs:
• Clear — conditions met, unambiguous
• Partial — some conditions unmet → defer to surrogate
• No coverage — directive silent; do NOT presume
• Vague — surface imprecise language (“no heroic measures”)

Safety feature: uncertainty is visible. A checkbox or overconfident LLM cannot say “I don’t know.”

Transition → end-to-end pipeline.""",

    # Slide 7
    """[~45s] LLMs ONLY at the boundary. Closed-world extraction → validated JSON → OWL individuals → deterministic scenario reasoner → four-category output + code status/POLST.

Closed-world rule: extract only what the patient stated; out-of-scope text (nutrition, antibiotics) must NOT map to nearest class—that enables trustworthy “no coverage.”

Tagline: ontology WITH LLMs, not versus LLMs.

Transition → how we evaluated.""",

    # Slide 8
    """[~50s] Five evaluation strands:
1) Pipeline: example patient populates end-to-end (9 prefs, 50 individuals).
2) 16 vignettes, all decision points & output types; gold = team-adjudicated expected answers.
3) Ablation (post-hoc sensitivity): re-score vignettes ignoring activation conditions → 12/16 (75%) vs 16/16 (100%). All four errors are conditional/partial cases—shows condition modeling matters. (Document as sensitivity analysis, not a separate SWRL experiment.)
4) LLM extraction: 12 real-template statements + 2 out-of-scope tests.
5) Track 3: 12 profiles → code status/POLST vs POLST semantics gold + Magnus expert review.

Caveat if asked: gold is team-authored but grounded in real templates and independent POLST definitions; blind rater sheets in repo.""",

    # Slide 9
    """[~60s — FLAGSHIP] Track 1: 16/16 decision AND match-type accuracy on vignettes.

Track 2: Precision 0.94, recall 1.00, F1 0.97; negation/type/conditionality 15/15; ZERO hallucinations on out-of-scope nutrition & antibiotics.

Track 3: 12 real template profiles → 35/36 field agreement (97%): code status 12/12, POLST A 12/12, POLST B 11/12, exact profile 11/12; κ = 1.00 / 1.00 / 0.88.

Ablation: ignoring conditions → 75% (4 false-confident errors on conditional vignettes).

If time: P9 vs P10 same directive flips DNR↔Full code; P5 “miss” is principled (Not specified vs presuming Full Treatment).

Transition → concrete if/then example.""",

    # Slide 10
    """[~50s] Same directive, two scenarios:
A: NYHA IV + no reversible cause → DNR, clear.
B: NYHA III + reversible cause → Full code, partial.

Flat POLST/checkbox must over-apply or discard the condition. ADO preserves if/then and flags when it applies.

Transition → where this lives clinically.""",

    # Slide 11
    """[~40s] Dr. David Magnus (Stanford Biomedical Ethics): “You can’t think of these as AHCDs… more like a progress note.”

Near-term product: pre-populate Epic-style ACP note → provider conversation can clarify/override → final ACP + POLST/code status.

Not a legal directive; complements POLST and code status.

If asked: next modeling gap is WHO decides (surrogate override), not just WHAT intervention.

Transition → limitations.""",

    # Slide 12
    """[~40s] Honest limits:
• No legal force
• One-dimensional (what, not who decides) — top priority extension
• Coverage gaps: artificial nutrition, antibiotics, comfort care
• Temporal: numeric hour bounds, not full temporal logic
• Gold standards team-authored; real EHR validation (MIMIC-IV) future

Principle: surface uncertainty; don’t replace clinicians/surrogates.""",

    # Slide 13
    """[~30s] Four takeaways on slide—read them, don’t rush.

Close: Questions welcome.
Repo: https://github.com/patrickwalsh26/bmds210-AD_Ontology

Q&A prep: legal AD? No. Why not LLM only? Over-confident, not auditable. Weakest point? Team gold + P5 principled divergence.""",
]


def set_notes(prs: Presentation) -> None:
    for i, slide in enumerate(prs.slides):
        if i >= len(NOTES):
            break
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.text = NOTES[i]


def update_slide_9(prs: Presentation) -> None:
    slide = prs.slides[8]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if t == "κ=1.00":
            shape.text_frame.text = "κ=1.00 / 1.00 / 0.88"
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
        elif t == "code status + POLST A":
            shape.text_frame.text = "code / POLST A / POLST B"
        elif "Ignoring activation conditions" in t:
            shape.text_frame.text = (
                "Post-hoc ablation: ignore activation conditions on the same 16 vignettes → 12/16 (75%). "
                "All four errors are conditional/partial cases."
            )

    # Footnote metrics box
    left, top, width, height = Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.55)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Also: extraction P/R = 0.94 / 1.00 • 0 out-of-scope hallucinations • "
        "11/12 profiles exact match on all three bedside fields"
    )
    p.font.size = Pt(11)


def main() -> None:
    prs = Presentation(str(PPTX))
    set_notes(prs)
    update_slide_9(prs)
    prs.save(str(PPTX))
    print(f"Updated {PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
